from __future__ import annotations

import io
import json
import math
import os
import re
import tempfile
import time
from dataclasses import dataclass, field
from typing import Any, Callable, Dict, Iterable, List, Optional, Sequence

import requests

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
    PPTX_IMPORT_ERROR = None
except ImportError as exc:  # pragma: no cover
    Presentation = None
    MSO_SHAPE_TYPE = None
    MSO_AUTO_SIZE = None
    MSO_VERTICAL_ANCHOR = None
    PP_ALIGN = None
    Inches = None
    Pt = None
    PPTX_IMPORT_ERROR = exc

try:
    import pythoncom
    from win32com.client import DispatchEx

    POWERPOINT_COM_ERROR = None
except ImportError as exc:  # pragma: no cover
    pythoncom = None
    DispatchEx = None
    POWERPOINT_COM_ERROR = exc


CHINESE_RE = re.compile(r"[\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff]")
NUMERIC_ONLY_RE = re.compile(r"^[\d\s.,:%+\-()/]+$")
POWERPOINT_PARAGRAPH_MARKS = "\r\x0b"
MSO_GROUP = 6
MSO_TRUE = -1
MSO_FALSE = 0
PPT_SAVE_AS_OPENXML = 24
PP_AUTO_SIZE_SHAPE_TO_FIT_TEXT = 1
PP_AUTO_SIZE_NONE = 0
MIN_FONT_SIZE_PT = 9.0
FONT_SHRINK_STEP = 0.92
SHAPE_EXPAND_GAP = 4.0
MAX_WIDTH_EXPAND_RATIO = 0.18
ALLOW_SAFE_WIDTH_EXPANSION = False
PP_ALIGN_LEFT = 1
PP_ALIGN_CENTER = 2
PP_ALIGN_RIGHT = 3
PP_ALIGN_JUSTIFY = 4
PP_ALIGN_DISTRIBUTE = 5
PP_ALIGN_THAI_DISTRIBUTE = 6
PP_ALIGN_JUSTIFY_LOW = 7
PP_VERTICAL_ANCHOR_TOP = 1


class PPTTranslationError(Exception):
    """Raised when the PPT translation flow cannot be completed safely."""


@dataclass
class TranslationUnit:
    unit_id: str
    text: str
    paragraph: Any
    shape: Any
    target_type: str
    paragraph_context: str = ""
    run: Optional[Any] = None
    suffix: str = ""


@dataclass
class ShapeLayoutContext:
    slide: Any
    shape: Any
    original_text: str
    original_height: float
    inside_group: bool = False


@dataclass
class TableLayoutContext:
    slide: Any
    shape: Any
    original_row_heights: List[float]
    original_height: float
    inside_group: bool = False


@dataclass
class SlideTranslationContext:
    units: List[TranslationUnit] = field(default_factory=list)
    shape_layouts: List[ShapeLayoutContext] = field(default_factory=list)
    table_layouts: List[TableLayoutContext] = field(default_factory=list)


def get_runtime_environment() -> Dict[str, Any]:
    if PPTX_IMPORT_ERROR is None:
        return {
            "ready": True,
            "engine": "python-pptx",
            "message": "当前使用 python-pptx 引擎，适合标准服务端处理。",
        }
    if POWERPOINT_COM_ERROR is None:
        return {
            "ready": True,
            "engine": "powerpoint-com",
            "message": (
                "当前未安装 python-pptx，已切换为本机 Microsoft PowerPoint 引擎。"
                " 这很适合你当前这台电脑自己使用。"
            ),
        }
    return {
        "ready": False,
        "engine": None,
        "message": (
            "当前环境既没有 python-pptx，也没有可用的 Microsoft PowerPoint COM 组件，"
            "暂时无法处理 PPT。"
        ),
    }


class GeminiTranslator:
    def __init__(
        self,
        api_key: str,
        model: str = "gemini-2.5-flash",
        timeout: int = 120,
        max_chars_per_batch: int = 3500,
    ) -> None:
        self.api_key = api_key.strip()
        self.model = model.strip() or "gemini-2.5-flash"
        self.timeout = timeout
        self.max_chars_per_batch = max_chars_per_batch
        if not self.api_key:
            raise PPTTranslationError("Gemini API Key 不能为空。")

    def translate_units(self, units: Sequence[TranslationUnit]) -> Dict[str, str]:
        results = {unit.unit_id: unit.text for unit in units}
        candidates = [unit for unit in units if should_translate(unit.text)]
        for batch in chunk_units(candidates, self.max_chars_per_batch):
            translated = self._translate_batch(batch)
            results.update(translated)
        return results

    def _translate_batch(self, batch: Sequence[TranslationUnit]) -> Dict[str, str]:
        payload_items = []
        for unit in batch:
            payload_items.append(
                {
                    "id": unit.unit_id,
                    "text": unit.text,
                    "context": unit.paragraph_context[:500],
                }
            )

        body = {
            "systemInstruction": {
                "parts": [
                    {
                        "text": (
                            "You are a professional presentation translator. "
                            "Translate Chinese into concise, business-ready English for PPT slides. "
                            "Preserve numbers, brand names, punctuation, and line breaks whenever possible. "
                            "Return JSON only."
                        )
                    }
                ]
            },
            "contents": [
                {
                    "parts": [
                        {
                            "text": json.dumps(
                                {
                                    "task": (
                                        "Translate every item into polished English suitable for a business PPT. "
                                        "If an item is already English, numeric, or should remain unchanged, "
                                        "return it exactly as-is. Keep the same ids."
                                    ),
                                    "output_schema": {
                                        "translations": [
                                            {"id": "string", "text": "string"},
                                        ]
                                    },
                                    "items": payload_items,
                                },
                                ensure_ascii=False,
                            )
                        }
                    ]
                }
            ],
            "generationConfig": {
                "temperature": 0.2,
                "responseMimeType": "application/json",
            },
        }

        url = "https://generativelanguage.googleapis.com/v1beta/models/{0}:generateContent".format(
            self.model
        )
        headers = {
            "Content-Type": "application/json",
            "x-goog-api-key": self.api_key,
        }

        last_error = None
        max_attempts = 3
        for attempt in range(max_attempts):
            try:
                response = requests.post(
                    url,
                    headers=headers,
                    json=body,
                    timeout=self.timeout,
                )
                if response.status_code == 429:
                    last_error = self._build_rate_limit_error(response)
                    if attempt < max_attempts - 1:
                        time.sleep(self._retry_delay_seconds(response, attempt))
                        continue
                    raise last_error
                response.raise_for_status()
                raw = response.json()
                content = extract_gemini_text(raw)
                return self._parse_translations(content, batch)
            except (KeyError, ValueError, requests.RequestException) as exc:
                last_error = exc

        message = "Gemini 翻译请求失败，请检查 API Key、模型名称、网络或接口额度。"
        if last_error is not None:
            message = "{0}\n原始错误: {1}".format(message, last_error)
        raise PPTTranslationError(message)

    def _build_rate_limit_error(self, response: requests.Response) -> PPTTranslationError:
        detail = ""
        try:
            payload = response.json()
            detail = safe_get(payload.get("error", {}), "message", "")
        except ValueError:
            detail = response.text[:300]

        message = (
            "Gemini 返回 429 Too Many Requests，当前请求被限流或额度不足。\n"
            "建议先切换为 `gemini-2.5-flash`，或稍后重试。"
        )
        if detail:
            message = "{0}\n接口信息: {1}".format(message, detail)
        return PPTTranslationError(message)

    def _retry_delay_seconds(self, response: requests.Response, attempt: int) -> float:
        retry_after = response.headers.get("Retry-After")
        if retry_after:
            try:
                return min(float(retry_after), 20.0)
            except ValueError:
                pass
        return min(2.0 * (attempt + 1), 6.0)

    def _parse_translations(
        self,
        content: str,
        batch: Sequence[TranslationUnit],
    ) -> Dict[str, str]:
        cleaned = strip_code_fence(content)
        data = json.loads(cleaned)
        translations = data.get("translations", [])
        lookup = {unit.unit_id: unit.text for unit in batch}
        for item in translations:
            unit_id = item.get("id")
            text = item.get("text")
            if unit_id in lookup and isinstance(text, str):
                lookup[unit_id] = text
        return lookup


def translate_presentation(
    pptx_bytes: bytes,
    api_key: str,
    progress_callback: Optional[Callable[[int, int, str], None]] = None,
    model: str = "gemini-2.5-flash",
) -> bytes:
    translator = GeminiTranslator(api_key=api_key, model=model)
    runtime = get_runtime_environment()
    if runtime["engine"] == "python-pptx":
        return translate_with_python_pptx(
            pptx_bytes=pptx_bytes,
            translator=translator,
            progress_callback=progress_callback,
        )
    if runtime["engine"] == "powerpoint-com":
        return translate_with_powerpoint_com(
            pptx_bytes=pptx_bytes,
            translator=translator,
            progress_callback=progress_callback,
        )
    raise PPTTranslationError(runtime["message"])


def translate_with_python_pptx(
    pptx_bytes: bytes,
    translator: GeminiTranslator,
    progress_callback: Optional[Callable[[int, int, str], None]] = None,
) -> bytes:
    presentation = Presentation(io.BytesIO(pptx_bytes))
    total_slides = len(presentation.slides)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        if progress_callback is not None:
            progress_callback(slide_index, total_slides, "正在扫描可翻译文本")

        context = collect_slide_context(slide)
        if not context.units:
            if progress_callback is not None:
                progress_callback(slide_index, total_slides, "本页未发现需要翻译的中文文本")
            continue

        if progress_callback is not None:
            progress_callback(
                slide_index,
                total_slides,
                "正在翻译 {0} 段文本".format(len(context.units)),
            )

        translations = translator.translate_units(context.units)

        if progress_callback is not None:
            progress_callback(slide_index, total_slides, "正在回写译文并微调排版")

        apply_translations(context.units, translations)
        adjust_slide_layout(context)

    if progress_callback is not None:
        progress_callback(total_slides, total_slides, "正在生成新的 PPT 文件")

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def translate_with_powerpoint_com(
    pptx_bytes: bytes,
    translator: GeminiTranslator,
    progress_callback: Optional[Callable[[int, int, str], None]] = None,
) -> bytes:
    if POWERPOINT_COM_ERROR is not None:
        raise PPTTranslationError("当前环境没有可用的 Microsoft PowerPoint COM 组件。")

    with tempfile.TemporaryDirectory() as temp_dir:
        input_path = os.path.join(temp_dir, "input.pptx")
        output_path = os.path.join(temp_dir, "translated_output.pptx")
        with open(input_path, "wb") as handle:
            handle.write(pptx_bytes)

        pythoncom.CoInitialize()
        application = None
        presentation = None
        try:
            application = DispatchEx("PowerPoint.Application")
            application.DisplayAlerts = 0
            presentation = application.Presentations.Open(
                input_path,
                False,
                False,
                False,
            )
            total_slides = int(presentation.Slides.Count)

            for slide_index in range(1, total_slides + 1):
                slide = presentation.Slides(slide_index)
                if progress_callback is not None:
                    progress_callback(slide_index, total_slides, "正在扫描可翻译文本")

                context = collect_com_slide_context(slide)
                if not context.units:
                    if progress_callback is not None:
                        progress_callback(slide_index, total_slides, "本页未发现需要翻译的中文文本")
                    continue

                if progress_callback is not None:
                    progress_callback(
                        slide_index,
                        total_slides,
                        "正在翻译 {0} 段文本".format(len(context.units)),
                    )

                translations = translator.translate_units(context.units)

                if progress_callback is not None:
                    progress_callback(slide_index, total_slides, "正在回写译文并微调排版")

                apply_com_translations(context.units, translations)
                adjust_com_slide_layout(context)

            if progress_callback is not None:
                progress_callback(total_slides, total_slides, "正在生成新的 PPT 文件")

            presentation.SaveAs(output_path, PPT_SAVE_AS_OPENXML)
            with open(output_path, "rb") as handle:
                return handle.read()
        except PPTTranslationError:
            raise
        except Exception as exc:
            raise PPTTranslationError(
                "PowerPoint 本机引擎处理失败，请确认本机 Office 可正常打开 PPT。\n原始错误: {0}".format(
                    exc
                )
            )
        finally:
            if presentation is not None:
                try:
                    presentation.Close()
                except Exception:
                    pass
            if application is not None:
                try:
                    application.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()


def collect_slide_context(slide: Any) -> SlideTranslationContext:
    context = SlideTranslationContext()
    counters = {"value": 0}

    def walk_shapes(shapes: Iterable[Any], inside_group: bool = False) -> None:
        for shape in shapes:
            shape_type = getattr(shape, "shape_type", None)
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                walk_shapes(shape.shapes, inside_group=True)
                continue

            if getattr(shape, "has_table", False):
                table_units = collect_table_units(shape, counters=counters)
                if table_units:
                    context.table_layouts.append(
                        TableLayoutContext(
                            slide=slide,
                            shape=shape,
                            original_row_heights=[float(row.height) for row in shape.table.rows],
                            original_height=float(shape.height),
                            inside_group=inside_group,
                        )
                    )
                    context.units.extend(table_units)
                continue

            if getattr(shape, "has_text_frame", False):
                text_units = collect_text_frame_units(
                    text_frame=shape.text_frame,
                    shape=shape,
                    counters=counters,
                )
                if text_units:
                    context.shape_layouts.append(
                        ShapeLayoutContext(
                            slide=slide,
                            shape=shape,
                            original_text=extract_text_from_text_frame(shape.text_frame),
                            original_height=float(shape.height),
                            inside_group=inside_group,
                        )
                    )
                    context.units.extend(text_units)

    walk_shapes(slide.shapes)
    return context


def collect_com_slide_context(slide: Any) -> SlideTranslationContext:
    context = SlideTranslationContext()
    counters = {"value": 0}

    def walk_shapes(shapes: Any, inside_group: bool = False) -> None:
        shape_count = int(shapes.Count)
        for index in range(1, shape_count + 1):
            shape = shapes.Item(index)
            if safe_get(shape, "Type") == MSO_GROUP:
                group_items = safe_get(shape, "GroupItems")
                if group_items is not None:
                    walk_shapes(group_items, inside_group=True)
                continue

            if com_shape_has_table(shape):
                table_units = collect_com_table_units(shape, counters=counters)
                if table_units:
                    table = shape.Table
                    context.table_layouts.append(
                        TableLayoutContext(
                            slide=slide,
                            shape=shape,
                            original_row_heights=[
                                float(table.Rows(row_index).Height)
                                for row_index in range(1, int(table.Rows.Count) + 1)
                            ],
                            original_height=float(shape.Height),
                            inside_group=inside_group,
                        )
                    )
                    context.units.extend(table_units)
                continue

            if com_shape_has_text(shape):
                text_units = collect_com_text_range_units(
                    text_range=shape.TextFrame.TextRange,
                    shape=shape,
                    counters=counters,
                )
                if text_units:
                    context.shape_layouts.append(
                        ShapeLayoutContext(
                            slide=slide,
                            shape=shape,
                            original_text=extract_text_from_com_text_range(shape.TextFrame.TextRange),
                            original_height=float(shape.Height),
                            inside_group=inside_group,
                        )
                    )
                    context.units.extend(text_units)

    walk_shapes(slide.Shapes)
    return context


def collect_table_units(shape: Any, counters: Dict[str, int]) -> List[TranslationUnit]:
    units: List[TranslationUnit] = []
    table = shape.table
    for row in table.rows:
        for cell in row.cells:
            units.extend(
                collect_text_frame_units(
                    text_frame=cell.text_frame,
                    shape=shape,
                    counters=counters,
                )
            )
    return units


def collect_com_table_units(shape: Any, counters: Dict[str, int]) -> List[TranslationUnit]:
    units: List[TranslationUnit] = []
    table = shape.Table
    for row_index in range(1, int(table.Rows.Count) + 1):
        for col_index in range(1, int(table.Columns.Count) + 1):
            cell_shape = table.Cell(row_index, col_index).Shape
            if not com_shape_has_text(cell_shape):
                continue
            units.extend(
                collect_com_text_range_units(
                    text_range=cell_shape.TextFrame.TextRange,
                    shape=shape,
                    counters=counters,
                )
            )
    return units


def collect_text_frame_units(
    text_frame: Any,
    shape: Any,
    counters: Dict[str, int],
) -> List[TranslationUnit]:
    units: List[TranslationUnit] = []
    for paragraph in text_frame.paragraphs:
        paragraph_text = paragraph.text or ""
        clean_text, suffix = split_suffix(paragraph_text)
        if not clean_text.strip():
            continue
        if not should_translate(clean_text):
            continue

        non_empty_runs = [run for run in paragraph.runs if run.text]
        mixed_styles = len(non_empty_runs) > 1 and not runs_share_style(non_empty_runs)

        if mixed_styles:
            for run in paragraph.runs:
                run_text = run.text or ""
                clean_run_text, run_suffix = split_suffix(run_text)
                if not clean_run_text or not should_translate(clean_run_text):
                    continue
                counters["value"] += 1
                units.append(
                    TranslationUnit(
                        unit_id="unit_{0}".format(counters["value"]),
                        text=clean_run_text,
                        paragraph=paragraph,
                        shape=shape,
                        target_type="run",
                        paragraph_context=clean_text,
                        run=run,
                        suffix=run_suffix,
                    )
                )
        else:
            counters["value"] += 1
            units.append(
                TranslationUnit(
                    unit_id="unit_{0}".format(counters["value"]),
                    text=clean_text,
                    paragraph=paragraph,
                    shape=shape,
                    target_type="paragraph",
                    paragraph_context=clean_text,
                    suffix=suffix,
                )
            )
    return units


def collect_com_text_range_units(
    text_range: Any,
    shape: Any,
    counters: Dict[str, int],
) -> List[TranslationUnit]:
    units: List[TranslationUnit] = []
    paragraph_count = int(text_range.Paragraphs().Count)
    for paragraph_index in range(1, paragraph_count + 1):
        paragraph = text_range.Paragraphs(paragraph_index, 1)
        paragraph_text = safe_text(paragraph)
        clean_text, suffix = split_suffix(paragraph_text)
        if not clean_text.strip():
            continue
        if not should_translate(clean_text):
            continue

        non_empty_runs = collect_non_empty_com_runs(paragraph)
        mixed_styles = len(non_empty_runs) > 1 and not com_runs_share_style(non_empty_runs)

        if mixed_styles:
            for run in non_empty_runs:
                run_text = safe_text(run)
                clean_run_text, run_suffix = split_suffix(run_text)
                if not clean_run_text or not should_translate(clean_run_text):
                    continue
                counters["value"] += 1
                units.append(
                    TranslationUnit(
                        unit_id="unit_{0}".format(counters["value"]),
                        text=clean_run_text,
                        paragraph=paragraph,
                        shape=shape,
                        target_type="run",
                        paragraph_context=clean_text,
                        run=run,
                        suffix=run_suffix,
                    )
                )
        else:
            counters["value"] += 1
            units.append(
                TranslationUnit(
                    unit_id="unit_{0}".format(counters["value"]),
                    text=clean_text,
                    paragraph=paragraph,
                    shape=shape,
                    target_type="paragraph",
                    paragraph_context=clean_text,
                    suffix=suffix,
                )
            )
    return units


def apply_translations(
    units: Sequence[TranslationUnit],
    translations: Dict[str, str],
) -> None:
    for unit in units:
        translated = normalize_translation_text(translations.get(unit.unit_id, unit.text))
        payload = translated + unit.suffix
        if unit.target_type == "run" and unit.run is not None:
            unit.run.text = payload
        else:
            replace_paragraph_text(unit.paragraph, payload)


def apply_com_translations(
    units: Sequence[TranslationUnit],
    translations: Dict[str, str],
) -> None:
    for unit in units:
        translated = normalize_translation_text(translations.get(unit.unit_id, unit.text))
        payload = translated + unit.suffix
        target = unit.run if unit.target_type == "run" and unit.run is not None else unit.paragraph
        try:
            target.Text = payload
        except Exception:
            raise PPTTranslationError("PowerPoint 文本回写失败，可能遇到受保护或异常文本框。")


def adjust_slide_layout(context: SlideTranslationContext) -> None:
    layout_items: List[Any] = list(context.shape_layouts) + list(context.table_layouts)
    layout_items.sort(key=lambda item: int(item.shape.top))
    for item in layout_items:
        if isinstance(item, ShapeLayoutContext):
            adjust_shape_layout(item)
        else:
            adjust_table_layout(item)


def adjust_com_slide_layout(context: SlideTranslationContext) -> None:
    layout_items: List[Any] = list(context.shape_layouts) + list(context.table_layouts)
    layout_items.sort(key=lambda item: float(safe_get(item.shape, "Top", 0.0)))
    for item in layout_items:
        if isinstance(item, ShapeLayoutContext):
            adjust_com_shape_layout(item)
        else:
            adjust_com_table_layout(item)


def adjust_shape_layout(layout: ShapeLayoutContext) -> int:
    shape = layout.shape
    if not getattr(shape, "has_text_frame", False):
        return 0

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    try:
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    except ValueError:
        pass
    shape.height = int(layout.original_height)
    normalize_python_text_frame_layout(text_frame)
    maybe_expand_python_shape_width(layout)

    translated_text = extract_text_from_text_frame(text_frame)
    if not translated_text.strip():
        return 0

    required_height = estimate_text_frame_height(text_frame, int(shape.width))
    current_height = int(shape.height)
    if required_height <= current_height:
        return 0

    shrink_text_frame_fonts_to_fit(text_frame, current_height, int(shape.width))
    return 0


def adjust_com_shape_layout(layout: ShapeLayoutContext) -> float:
    shape = layout.shape
    if not com_shape_has_text(shape):
        return 0.0

    old_height = float(layout.original_height)
    shape.Height = old_height
    text_frame = shape.TextFrame
    try:
        text_frame.WordWrap = MSO_TRUE
        text_frame.AutoSize = PP_AUTO_SIZE_NONE
    except Exception:
        pass
    normalize_com_text_frame_layout(text_frame)
    maybe_expand_com_shape_width(layout)

    translated_text = extract_text_from_com_text_range(text_frame.TextRange)
    if not translated_text.strip():
        return 0.0

    required_height = estimate_com_text_range_height(text_frame.TextRange, float(shape.Width), text_frame)
    if required_height <= old_height:
        return 0.0

    shrink_com_text_range_fonts_to_fit(text_frame.TextRange, text_frame, old_height, float(shape.Width))
    return 0.0


def adjust_table_layout(layout: TableLayoutContext) -> int:
    table = layout.shape.table
    for row_index, row in enumerate(table.rows):
        row.height = int(layout.original_row_heights[row_index])
        current_height = int(row.height)
        for col_index, cell in enumerate(row.cells):
            cell.text_frame.word_wrap = True
            try:
                cell.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            except ValueError:
                pass
            normalize_python_text_frame_layout(cell.text_frame)
            cell_width = int(table.columns[col_index].width)
            estimated = estimate_text_frame_height(cell.text_frame, cell_width)
            if estimated > current_height:
                shrink_text_frame_fonts_to_fit(cell.text_frame, current_height, cell_width)
    layout.shape.height = int(layout.original_height)
    return 0


def adjust_com_table_layout(layout: TableLayoutContext) -> float:
    table = layout.shape.Table
    row_count = int(table.Rows.Count)
    col_count = int(table.Columns.Count)
    for row_index in range(1, row_count + 1):
        row = table.Rows(row_index)
        original_height = float(layout.original_row_heights[row_index - 1])
        row.Height = original_height
        for col_index in range(1, col_count + 1):
            cell_shape = table.Cell(row_index, col_index).Shape
            if not com_shape_has_text(cell_shape):
                continue
            try:
                cell_shape.TextFrame.WordWrap = MSO_TRUE
                cell_shape.TextFrame.AutoSize = PP_AUTO_SIZE_NONE
            except Exception:
                pass
            normalize_com_text_frame_layout(cell_shape.TextFrame)
            estimated = estimate_com_text_range_height(
                cell_shape.TextFrame.TextRange,
                float(cell_shape.Width),
                cell_shape.TextFrame,
            )
            if estimated > original_height:
                shrink_com_text_range_fonts_to_fit(
                    cell_shape.TextFrame.TextRange,
                    cell_shape.TextFrame,
                    original_height,
                    float(cell_shape.Width),
                )
    layout.shape.Height = float(layout.original_height)
    return 0.0


def shrink_text_frame_fonts_to_fit(text_frame: Any, max_height_emu: int, width_emu: int) -> None:
    targets = collect_python_font_targets(text_frame)
    if not targets:
        return

    for _ in range(12):
        required_height = estimate_text_frame_height(text_frame, width_emu)
        if required_height <= max_height_emu:
            return
        if min(target["size_pt"] for target in targets) <= MIN_FONT_SIZE_PT:
            return
        for target in targets:
            next_size = max(MIN_FONT_SIZE_PT, target["size_pt"] * FONT_SHRINK_STEP)
            set_python_font_target_size(target, next_size)


def shrink_com_text_range_fonts_to_fit(
    text_range: Any,
    text_frame: Any,
    max_height_points: float,
    width_points: float,
) -> None:
    targets = collect_com_font_targets(text_range)
    if not targets:
        return

    for _ in range(12):
        required_height = estimate_com_text_range_height(text_range, width_points, text_frame)
        if required_height <= max_height_points:
            return
        if min(target["size_pt"] for target in targets) <= MIN_FONT_SIZE_PT:
            return
        for target in targets:
            next_size = max(MIN_FONT_SIZE_PT, target["size_pt"] * FONT_SHRINK_STEP)
            set_com_font_target_size(target, next_size)


def normalize_python_text_frame_layout(text_frame: Any) -> None:
    try:
        if MSO_VERTICAL_ANCHOR is not None:
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    except Exception:
        pass
    for paragraph in text_frame.paragraphs:
        text = paragraph.text or ""
        if not text.strip():
            continue
        if contains_latin_text(text) and PP_ALIGN is not None and paragraph.alignment in {
            PP_ALIGN.JUSTIFY,
            PP_ALIGN.JUSTIFY_LOW,
            PP_ALIGN.DISTRIBUTE,
            PP_ALIGN.THAI_DISTRIBUTE,
        }:
            paragraph.alignment = PP_ALIGN.LEFT
        try:
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            paragraph.line_spacing = 1.0
        except Exception:
            pass


def normalize_com_text_frame_layout(text_frame: Any) -> None:
    try:
        text_frame.VerticalAnchor = PP_VERTICAL_ANCHOR_TOP
    except Exception:
        pass
    text_range = text_frame.TextRange
    paragraph_count = int(text_range.Paragraphs().Count)
    for paragraph_index in range(1, paragraph_count + 1):
        paragraph = text_range.Paragraphs(paragraph_index, 1)
        text = safe_text(paragraph)
        if not strip_paragraph_marks(text).strip():
            continue
        paragraph_format = safe_get(paragraph, "ParagraphFormat")
        alignment = safe_get(paragraph_format, "Alignment")
        if contains_latin_text(text) and alignment in {
            PP_ALIGN_JUSTIFY,
            PP_ALIGN_JUSTIFY_LOW,
            PP_ALIGN_DISTRIBUTE,
            PP_ALIGN_THAI_DISTRIBUTE,
        }:
            paragraph_format.Alignment = PP_ALIGN_LEFT
        try:
            paragraph_format.LineRuleBefore = MSO_FALSE
            paragraph_format.SpaceBefore = 0
            paragraph_format.LineRuleAfter = MSO_FALSE
            paragraph_format.SpaceAfter = 0
            paragraph_format.LineRuleWithin = MSO_TRUE
            paragraph_format.SpaceWithin = 1.0
        except Exception:
            pass


def maybe_expand_python_shape_width(layout: ShapeLayoutContext) -> None:
    if not ALLOW_SAFE_WIDTH_EXPANSION:
        return
    shape = layout.shape
    slide = layout.slide
    try:
        original_width = int(shape.width)
        target_width = compute_python_safe_expanded_width(slide, shape, original_width)
        if target_width > original_width:
            shape.width = target_width
    except Exception:
        pass


def maybe_expand_com_shape_width(layout: ShapeLayoutContext) -> None:
    if not ALLOW_SAFE_WIDTH_EXPANSION:
        return
    shape = layout.shape
    slide = layout.slide
    try:
        original_width = float(shape.Width)
        target_width = compute_com_safe_expanded_width(slide, shape, original_width)
        if target_width > original_width:
            shape.Width = target_width
    except Exception:
        pass


def compute_python_safe_expanded_width(slide: Any, anchor_shape: Any, original_width: int) -> int:
    left = int(anchor_shape.left)
    top = int(anchor_shape.top)
    bottom = int(anchor_shape.top + anchor_shape.height)
    max_right = left + int(original_width * (1.0 + MAX_WIDTH_EXPAND_RATIO))
    safe_right = max_right

    try:
        presentation = slide.part.package.presentation_part.presentation
        slide_width = int(presentation.slide_width)
        safe_right = min(safe_right, slide_width - int(Inches(0.03)))
    except Exception:
        pass

    for shape in slide.shapes:
        if shape == anchor_shape:
            continue
        try:
            other_left = int(shape.left)
            other_top = int(shape.top)
            other_bottom = int(shape.top + shape.height)
        except Exception:
            continue
        if min(bottom, other_bottom) - max(top, other_top) <= 0:
            continue
        if other_left <= left + original_width:
            continue
        safe_right = min(safe_right, other_left - int(Inches(0.03)))

    return max(original_width, safe_right - left)


def compute_com_safe_expanded_width(slide: Any, anchor_shape: Any, original_width: float) -> float:
    left = float(anchor_shape.Left)
    top = float(anchor_shape.Top)
    bottom = float(anchor_shape.Top + anchor_shape.Height)
    max_right = left + original_width * (1.0 + MAX_WIDTH_EXPAND_RATIO)
    safe_right = max_right

    slide_width = float(
        safe_get(
            safe_get(safe_get(slide, "Parent"), "PageSetup"),
            "SlideWidth",
            max_right,
        )
    )
    if slide_width > 0:
        safe_right = min(safe_right, slide_width - SHAPE_EXPAND_GAP)

    shape_count = int(slide.Shapes.Count)
    for index in range(1, shape_count + 1):
        shape = slide.Shapes(index)
        if safe_get(shape, "Id") == safe_get(anchor_shape, "Id"):
            continue
        other_left = float(safe_get(shape, "Left", 0.0))
        other_top = float(safe_get(shape, "Top", 0.0))
        other_bottom = float(safe_get(shape, "Top", 0.0) + safe_get(shape, "Height", 0.0))
        if min(bottom, other_bottom) - max(top, other_top) <= 0:
            continue
        if other_left <= left + original_width:
            continue
        safe_right = min(safe_right, other_left - SHAPE_EXPAND_GAP)

    return max(original_width, safe_right - left)


def collect_python_font_targets(text_frame: Any) -> List[Dict[str, Any]]:
    targets: List[Dict[str, Any]] = []
    for paragraph in text_frame.paragraphs:
        non_empty_runs = [run for run in paragraph.runs if (run.text or "").strip()]
        if non_empty_runs:
            for run in non_empty_runs:
                size_pt = average_font_size(paragraph)
                if run.font.size is not None:
                    size_pt = float(run.font.size.pt)
                targets.append({"kind": "run", "obj": run, "size_pt": size_pt})
        else:
            targets.append(
                {
                    "kind": "paragraph",
                    "obj": paragraph,
                    "size_pt": average_font_size(paragraph),
                }
            )
    return targets


def set_python_font_target_size(target: Dict[str, Any], size_pt: float) -> None:
    target["size_pt"] = size_pt
    if target["kind"] == "run":
        target["obj"].font.size = Pt(size_pt)
    else:
        target["obj"].font.size = Pt(size_pt)


def collect_com_font_targets(text_range: Any) -> List[Dict[str, Any]]:
    targets: List[Dict[str, Any]] = []
    paragraph_count = int(text_range.Paragraphs().Count)
    for paragraph_index in range(1, paragraph_count + 1):
        paragraph = text_range.Paragraphs(paragraph_index, 1)
        non_empty_runs = collect_non_empty_com_runs(paragraph)
        if non_empty_runs:
            for run in non_empty_runs:
                size_pt = average_com_font_size(paragraph)
                run_size = safe_get(safe_get(run, "Font"), "Size")
                if run_size:
                    size_pt = float(run_size)
                targets.append({"kind": "run", "obj": run, "size_pt": size_pt})
        else:
            targets.append(
                {
                    "kind": "paragraph",
                    "obj": paragraph,
                    "size_pt": average_com_font_size(paragraph),
                }
            )
    return targets


def set_com_font_target_size(target: Dict[str, Any], size_pt: float) -> None:
    target["size_pt"] = size_pt
    font = safe_get(target["obj"], "Font")
    if font is not None:
        font.Size = size_pt


def shift_overlapping_shapes_down(slide: Any, anchor_shape: Any, delta: int) -> None:
    left = int(anchor_shape.left)
    right = int(anchor_shape.left + anchor_shape.width)
    new_bottom = int(anchor_shape.top + anchor_shape.height)
    gap = int(Inches(0.05))

    for shape in slide.shapes:
        if shape == anchor_shape:
            continue
        if not hasattr(shape, "left") or not hasattr(shape, "top"):
            continue

        other_left = int(shape.left)
        other_right = int(shape.left + shape.width)
        horizontal_overlap = min(right, other_right) - max(left, other_left)
        if horizontal_overlap <= 0:
            continue

        if int(shape.top) >= new_bottom + gap:
            continue
        if int(shape.top) < int(anchor_shape.top):
            continue

        shift_by = (new_bottom + gap) - int(shape.top)
        if shift_by > 0:
            shape.top = int(shape.top + shift_by)


def shift_overlapping_com_shapes_down(slide: Any, anchor_shape: Any, delta: float) -> None:
    left = float(anchor_shape.Left)
    right = float(anchor_shape.Left + anchor_shape.Width)
    new_bottom = float(anchor_shape.Top + anchor_shape.Height)
    gap = 3.5
    shape_count = int(slide.Shapes.Count)

    for index in range(1, shape_count + 1):
        shape = slide.Shapes(index)
        if safe_get(shape, "Id") == safe_get(anchor_shape, "Id"):
            continue

        other_left = float(safe_get(shape, "Left", 0.0))
        other_right = float(safe_get(shape, "Left", 0.0) + safe_get(shape, "Width", 0.0))
        horizontal_overlap = min(right, other_right) - max(left, other_left)
        if horizontal_overlap <= 0:
            continue

        other_top = float(safe_get(shape, "Top", 0.0))
        if other_top >= new_bottom + gap:
            continue
        if other_top < float(anchor_shape.Top):
            continue

        shift_by = (new_bottom + gap) - other_top
        if shift_by > 0:
            shape.Top = other_top + shift_by


def estimate_text_frame_height(text_frame: Any, width_emu: int) -> int:
    usable_width = max(
        int(Inches(0.6)),
        width_emu - int(text_frame.margin_left) - int(text_frame.margin_right),
    )
    total_height_pt = 0.0
    for paragraph in text_frame.paragraphs:
        text = paragraph.text or ""
        font_pt = average_font_size(paragraph)
        line_height_pt = max(font_pt * 1.25, 14.0)
        chars_per_line = max(6, int(emu_to_points(usable_width) / max(font_pt * 0.55, 1)))
        lines = max(1, int(math.ceil(visual_length(text) / float(chars_per_line))))
        total_height_pt += lines * line_height_pt

    total_height_pt += emu_to_points(int(text_frame.margin_top) + int(text_frame.margin_bottom))
    total_height_pt += 4.0
    return int(Pt(total_height_pt))


def estimate_com_text_range_height(text_range: Any, width_points: float, text_frame: Any) -> float:
    margin_left = float(safe_get(text_frame, "MarginLeft", 0.0))
    margin_right = float(safe_get(text_frame, "MarginRight", 0.0))
    margin_top = float(safe_get(text_frame, "MarginTop", 0.0))
    margin_bottom = float(safe_get(text_frame, "MarginBottom", 0.0))

    measured_height = measure_com_text_range_height(text_range, margin_top, margin_bottom)
    if measured_height > 0:
        return measured_height

    usable_width = max(40.0, width_points - margin_left - margin_right)

    total_height = 0.0
    paragraph_count = int(text_range.Paragraphs().Count)
    for paragraph_index in range(1, paragraph_count + 1):
        paragraph = text_range.Paragraphs(paragraph_index, 1)
        text = strip_paragraph_marks(safe_text(paragraph))
        font_pt = average_com_font_size(paragraph)
        paragraph_format = safe_get(paragraph, "ParagraphFormat")
        left_margin = max(0.0, float(safe_get(paragraph_format, "LeftMargin", 0.0)))
        first_margin = max(0.0, float(safe_get(paragraph_format, "FirstMargin", 0.0)))
        paragraph_width = max(24.0, usable_width - left_margin - first_margin)
        line_spacing = float(safe_get(paragraph_format, "SpaceWithin", 1.0) or 1.0)
        line_height_pt = max(font_pt * 1.15 * max(line_spacing, 1.0), 12.0)
        chars_per_line = max(4, int(paragraph_width / max(font_pt * 0.55, 1)))
        lines = max(1, int(math.ceil(visual_length(text) / float(chars_per_line))))
        total_height += lines * line_height_pt

    return total_height + margin_top + margin_bottom + 4.0


def measure_com_text_range_height(
    text_range: Any,
    margin_top: float,
    margin_bottom: float,
) -> float:
    try:
        bound_height = float(safe_get(text_range, "BoundHeight", 0.0) or 0.0)
    except Exception:
        return 0.0
    if bound_height <= 0:
        return 0.0
    return bound_height + margin_top + margin_bottom + 2.0


def average_font_size(paragraph: Any) -> float:
    sizes = []
    for run in paragraph.runs:
        if run.font.size is not None:
            sizes.append(float(run.font.size.pt))
    if sizes:
        return sum(sizes) / float(len(sizes))
    if paragraph.font.size is not None:
        return float(paragraph.font.size.pt)
    return 18.0


def average_com_font_size(paragraph: Any) -> float:
    sizes = []
    for run in collect_non_empty_com_runs(paragraph):
        size = safe_get(safe_get(run, "Font"), "Size")
        if size:
            sizes.append(float(size))
    if sizes:
        return sum(sizes) / float(len(sizes))
    font_size = safe_get(safe_get(paragraph, "Font"), "Size")
    if font_size:
        return float(font_size)
    return 18.0


def runs_share_style(runs: Sequence[Any]) -> bool:
    signatures = {run_style_signature(run) for run in runs if run.text}
    return len(signatures) <= 1


def com_runs_share_style(runs: Sequence[Any]) -> bool:
    signatures = {com_run_style_signature(run) for run in runs if strip_paragraph_marks(safe_text(run))}
    return len(signatures) <= 1


def run_style_signature(run: Any) -> Any:
    font = run.font
    color_rgb = None
    try:
        if font.color is not None and font.color.rgb is not None:
            color_rgb = str(font.color.rgb)
    except AttributeError:
        color_rgb = None
    return (
        font.name,
        float(font.size.pt) if font.size is not None else None,
        font.bold,
        font.italic,
        font.underline,
        color_rgb,
    )


def com_run_style_signature(run: Any) -> Any:
    font = safe_get(run, "Font")
    color_rgb = safe_get(safe_get(font, "Color"), "RGB")
    return (
        safe_get(font, "Name"),
        safe_get(font, "Size"),
        safe_get(font, "Bold"),
        safe_get(font, "Italic"),
        safe_get(font, "Underline"),
        color_rgb,
    )


def replace_paragraph_text(paragraph: Any, text: str) -> None:
    if paragraph.runs:
        first_run = paragraph.runs[0]
        first_run.text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run().text = text


def extract_text_from_text_frame(text_frame: Any) -> str:
    texts = [paragraph.text for paragraph in text_frame.paragraphs if paragraph.text]
    return "\n".join(texts)


def extract_text_from_com_text_range(text_range: Any) -> str:
    return strip_paragraph_marks(safe_text(text_range))


def should_translate(text: str) -> bool:
    stripped = (text or "").strip()
    if not stripped:
        return False
    if NUMERIC_ONLY_RE.match(stripped):
        return False
    return bool(CHINESE_RE.search(stripped))


def normalize_translation_text(text: str) -> str:
    return (text or "").replace("\r\n", "\n").replace("\u00a0", " ")


def looks_like_spaced_english(text: str) -> bool:
    cleaned = strip_paragraph_marks(text or "")
    return bool(re.search(r"[A-Za-z]{2,}", cleaned) and " " in cleaned)


def contains_latin_text(text: str) -> bool:
    cleaned = strip_paragraph_marks(text or "")
    return bool(re.search(r"[A-Za-z]{2,}", cleaned))


def strip_code_fence(text: str) -> str:
    cleaned = (text or "").strip()
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```[a-zA-Z0-9_-]*\n", "", cleaned)
        cleaned = re.sub(r"\n```$", "", cleaned)
    return cleaned


def chunk_units(
    units: Sequence[TranslationUnit],
    max_chars_per_batch: int,
) -> List[List[TranslationUnit]]:
    batches: List[List[TranslationUnit]] = []
    current: List[TranslationUnit] = []
    current_size = 0

    for unit in units:
        estimated = len(unit.text) + len(unit.paragraph_context)
        if current and current_size + estimated > max_chars_per_batch:
            batches.append(current)
            current = []
            current_size = 0
        current.append(unit)
        current_size += estimated

    if current:
        batches.append(current)
    return batches


def visual_length(text: str) -> float:
    score = 0.0
    for char in text or "":
        if char.isspace():
            score += 0.35
        elif CHINESE_RE.search(char):
            score += 1.0
        elif char.isupper():
            score += 0.72
        else:
            score += 0.58
    return score


def emu_to_points(value: int) -> float:
    return float(value) / 12700.0


def extract_gemini_text(data: Dict[str, Any]) -> str:
    if "error" in data:
        message = safe_get(data["error"], "message", "Gemini API 返回错误。")
        raise PPTTranslationError(str(message))

    candidates = data.get("candidates", [])
    for candidate in candidates:
        parts = safe_get(safe_get(candidate, "content", {}), "parts", [])
        if not isinstance(parts, list):
            continue
        text = "".join(part.get("text", "") for part in parts if isinstance(part, dict))
        if text:
            return text
    raise PPTTranslationError("Gemini 没有返回可解析的文本结果。")


def safe_get(obj: Any, attr: str, default: Any = None) -> Any:
    try:
        if isinstance(obj, dict):
            return obj.get(attr, default)
        return getattr(obj, attr)
    except Exception:
        return default


def safe_text(text_range: Any) -> str:
    value = safe_get(text_range, "Text", "")
    return value or ""


def strip_paragraph_marks(text: str) -> str:
    return (text or "").rstrip(POWERPOINT_PARAGRAPH_MARKS)


def split_suffix(text: str) -> Any:
    raw = text or ""
    core = raw.rstrip(POWERPOINT_PARAGRAPH_MARKS)
    suffix = raw[len(core):]
    return core, suffix


def collect_non_empty_com_runs(paragraph: Any) -> List[Any]:
    runs: List[Any] = []
    runs_collection = paragraph.Runs()
    run_count = int(runs_collection.Count)
    for run_index in range(1, run_count + 1):
        run = paragraph.Runs(run_index, 1)
        if strip_paragraph_marks(safe_text(run)):
            runs.append(run)
    return runs


def com_shape_has_table(shape: Any) -> bool:
    return safe_get(shape, "HasTable", 0) == MSO_TRUE


def com_shape_has_text(shape: Any) -> bool:
    try:
        return shape.HasTextFrame == MSO_TRUE and shape.TextFrame.HasText == MSO_TRUE
    except Exception:
        return False
